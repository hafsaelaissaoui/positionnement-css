from pptx import Presentation # Create a presentation object prs =
Presentation() # Slide 1: Title Slide slide_1 =
prs.slides.add_slide(prs.slide_layouts[0]) title_1 = slide_1.shapes.title
subtitle_1 = slide_1.placeholders[1] title_1.text = "Rapport de Stage chez TE
Connectivity" subtitle_1.text = "[Votre Nom] – [Date]" # Slide 2: Introduction
slide_2 = prs.slides.add_slide(prs.slide_layouts[1]) title_2 =
slide_2.shapes.title title_2.text = "Introduction" content_2 =
slide_2.shapes.placeholders[1].text_frame content_2.text = ("J'ai effectué mon
stage chez TE Connectivity, une entreprise reconnue dans le domaine de
l'automobile." " TE Connectivity emploie environ 80 000 personnes à travers le
monde, avec une présence nationale et internationale." " L'entreprise est un
leader dans la fabrication de connecteurs et de solutions de connectivité,
jouant un rôle essentiel" " dans l'automobile, l'aérospatial et d'autres
industries.") # Slide 3: Activités Principales de l'Entreprise slide_3 =
prs.slides.add_slide(prs.slide_layouts[1]) title_3 = slide_3.shapes.title
title_3.text = "Activités Principales de l'Entreprise" content_3 =
slide_3.shapes.placeholders[1].text_frame content_3.text = ("TE Connectivity
propose une large gamme de produits, notamment des connecteurs, des capteurs et
des solutions de communication," " spécifiquement pour le secteur automobile.
L'entreprise investit dans des technologies avancées comme l'automatisation" "
et le développement durable.") # Slide 4: Le Service de Maintenance slide_4 =
prs.slides.add_slide(prs.slide_layouts[1]) title_4 = slide_4.shapes.title
title_4.text = "Le Service de Maintenance" content_4 =
slide_4.shapes.placeholders[1].text_frame content_4.text = ("J'ai été affecté au
service de maintenance, qui est crucial pour assurer le bon fonctionnement des
équipements de production." " L'objectif principal de ce service est de garantir
la disponibilité et l'efficacité des machines, en prévenant les pannes.") #
Slide 5: Tâches Réalisées slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title_5 = slide_5.shapes.title title_5.text = "Tâches Réalisées" content_5 =
slide_5.shapes.placeholders[1].text_frame content_5.text = ("Au cours de mon
stage, j'ai participé à la maintenance préventive et corrective des machines." "
J'ai effectué des inspections régulières et des réparations sur des équipements
critiques." " J'ai également appris à utiliser des outils de diagnostic et de
réparation, ce qui m'a permis de mieux comprendre le fonctionnement des
machines.") # Slide 6: Collaboration Interdépartementale slide_6 =
prs.slides.add_slide(prs.slide_layouts[1]) title_6 = slide_6.shapes.title
title_6.text = "Collaboration Interdépartementale" content_6 =
slide_6.shapes.placeholders[1].text_frame content_6.text = ("J'ai travaillé en
étroite collaboration avec le service de production pour comprendre les besoins
en maintenance et optimiser les temps d'arrêt." " Cette collaboration a permis
d'assurer la disponibilité continue des équipements.") # Slide 7: Résultats
Obtenus slide_7 = prs.slides.add_slide(prs.slide_layouts[1]) title_7 =
slide_7.shapes.title title_7.text = "Résultats Obtenus" content_7 =
slide_7.shapes.placeholders[1].text_frame content_7.text = ("Mes contributions
ont permis de réduire les temps d'arrêt des machines, ce qui a conduit à une
augmentation de l'efficacité de la production.") # Slide 8: Apports du Stage
slide_8 = prs.slides.add_slide(prs.slide_layouts[1]) title_8 =
slide_8.shapes.title title_8.text = "Apports du Stage" content_8 =
slide_8.shapes.placeholders[1].text_frame content_8.text = ("J'ai acquis des
compétences en diagnostic de pannes et en maintenance préventive, ainsi que la
maîtrise de nouveaux outils." " Ce stage m'a également permis de développer mes
compétences en communication et en travail d'équipe, tout en découvrant les
enjeux de l'industrie automobile.") # Slide 9: Impact sur mon Projet
Professionnel slide_9 = prs.slides.add_slide(prs.slide_layouts[1]) title_9 =
slide_9.shapes.title title_9.text = "Impact sur mon Projet Professionnel"
content_9 = slide_9.shapes.placeholders[1].text_frame content_9.text = ("Ce
stage m'a aidé à mieux définir mes objectifs professionnels et à confirmer mon
intérêt pour le domaine de la maintenance industrielle.") # Slide 10: Conclusion
slide_10 = prs.slides.add_slide(prs.slide_layouts[1]) title_10 =
slide_10.shapes.title title_10.text = "Conclusion" content_10 =
slide_10.shapes.placeholders[1].text_frame content_10.text = ("Je tiens à
remercier TE Connectivity pour cette opportunité et tous mes collègues pour leur
soutien." " Ce stage a été très enrichissant, tant sur le plan professionnel que
personnel.") # Save presentation ppt_path =
"/mnt/data/Rapport_Stage_TE_Connectivity.pptx" prs.save(ppt_path) ppt_path
